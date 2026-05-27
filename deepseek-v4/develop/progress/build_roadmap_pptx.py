"""Generate the DeepSeek-V4 in Primus roadmap PPT (tech-style, black background).

Output: deepseek-v4/develop/progress/deepseek_v4_roadmap.pptx
"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUT_PATH = os.path.join(OUT_DIR, "deepseek_v4_roadmap.pptx")

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# AMD-inspired tech palette (deep black + neon accents)
BG_BLACK = RGBColor(0x05, 0x06, 0x0A)
BG_DEEP = RGBColor(0x0A, 0x0D, 0x18)
BG_PANEL = RGBColor(0x11, 0x16, 0x2A)
BG_PANEL_2 = RGBColor(0x18, 0x1F, 0x3A)
TEXT = RGBColor(0xE3, 0xE8, 0xFA)
TEXT_DIM = RGBColor(0x96, 0xA1, 0xC4)
TEXT_MUTE = RGBColor(0x5B, 0x66, 0x88)
LINE = RGBColor(0x33, 0x40, 0x6E)

ACCENT_CYAN = RGBColor(0x4D, 0xDC, 0xFF)
ACCENT_VIOLET = RGBColor(0xA8, 0x6B, 0xFF)
ACCENT_AMBER = RGBColor(0xFF, 0xB4, 0x6B)
ACCENT_EMERALD = RGBColor(0x5F, 0xFF, 0xA6)
ACCENT_ROSE = RGBColor(0xFF, 0x6B, 0x9B)
ACCENT_RED = RGBColor(0xED, 0x22, 0x4D)


# ----------------------------- helpers --------------------------------------


def set_slide_bg(slide, color: RGBColor = BG_BLACK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_rect(
    slide,
    x,
    y,
    w,
    h,
    fill=BG_PANEL,
    line=None,
    line_color=LINE,
    line_w=0.75,
    radius=None,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, x, y, w, h)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is False:
        s.line.fill.background()
    else:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_w)
    if radius is not None:
        # adjust corner radius (0..0.5)
        s.adjustments[0] = radius
    return s


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=14,
    color=TEXT,
    bold=False,
    italic=False,
    font="Inter",
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    spacing=1.15,
):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        lines = [text]
    else:
        lines = text
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        # Allow per-paragraph dict override
        if isinstance(ln, dict):
            p_text = ln.get("text", "")
            p_size = ln.get("size", size)
            p_color = ln.get("color", color)
            p_bold = ln.get("bold", bold)
            p_italic = ln.get("italic", italic)
            p_font = ln.get("font", font)
        else:
            p_text = ln
            p_size = size
            p_color = color
            p_bold = bold
            p_italic = italic
            p_font = font
        run = p.add_run()
        run.text = p_text
        run.font.name = p_font
        run.font.size = Pt(p_size)
        run.font.bold = p_bold
        run.font.italic = p_italic
        run.font.color.rgb = p_color
    return tb


def add_line(slide, x1, y1, x2, y2, color=LINE, weight=0.75):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    line.shadow.inherit = False
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def header_strip(slide, eyebrow, title, subtitle=None):
    """Top strip: small eyebrow text, big title, thin gradient bar."""
    add_text(
        slide,
        eyebrow.upper(),
        Inches(0.55),
        Inches(0.45),
        Inches(8),
        Inches(0.3),
        size=10,
        color=ACCENT_CYAN,
        bold=True,
        font="JetBrains Mono",
    )
    add_text(
        slide,
        title,
        Inches(0.55),
        Inches(0.75),
        Inches(12.2),
        Inches(0.7),
        size=30,
        color=TEXT,
        bold=True,
        font="Inter",
    )
    if subtitle:
        add_text(
            slide,
            subtitle,
            Inches(0.55),
            Inches(1.45),
            Inches(12.2),
            Inches(0.4),
            size=12,
            color=TEXT_DIM,
            font="Inter",
        )
    # accent bar
    bar = add_rect(
        slide,
        Inches(0.55),
        Inches(1.85),
        Inches(0.6),
        Emu(28575),
        fill=ACCENT_CYAN,
        line=False,
    )
    bar2 = add_rect(
        slide,
        Inches(1.2),
        Inches(1.85),
        Inches(0.3),
        Emu(28575),
        fill=ACCENT_VIOLET,
        line=False,
    )
    bar3 = add_rect(
        slide,
        Inches(1.55),
        Inches(1.85),
        Inches(0.15),
        Emu(28575),
        fill=ACCENT_EMERALD,
        line=False,
    )


def footer(slide, page_no, total):
    add_text(
        slide,
        "DeepSeek-V4 × Primus  ·  Plan-2 Roadmap",
        Inches(0.55),
        Inches(7.05),
        Inches(8),
        Inches(0.3),
        size=9,
        color=TEXT_MUTE,
        font="JetBrains Mono",
    )
    add_text(
        slide,
        f"{page_no:02d} / {total:02d}",
        Inches(11.8),
        Inches(7.05),
        Inches(1.0),
        Inches(0.3),
        size=9,
        color=TEXT_MUTE,
        font="JetBrains Mono",
        align=PP_ALIGN.RIGHT,
    )


def grid_lines(slide, rows=12, cols=20, color=None):
    """Subtle grid pattern in the background."""
    if color is None:
        color = RGBColor(0x12, 0x18, 0x28)
    # vertical
    for i in range(1, cols):
        x = SLIDE_W * i / cols
        add_line(slide, x, 0, x, SLIDE_H, color=color, weight=0.25)
    for i in range(1, rows):
        y = SLIDE_H * i / rows
        add_line(slide, 0, y, SLIDE_W, y, color=color, weight=0.25)


def chip(slide, x, y, label, color=ACCENT_CYAN, w=None):
    if w is None:
        w = Inches(1.0)
    h = Inches(0.28)
    s = add_rect(slide, x, y, w, h, fill=None, line_color=color, line_w=1.0, radius=0.5)
    add_text(
        slide,
        label,
        x,
        y + Emu(2000),
        w,
        h,
        size=9,
        color=color,
        bold=True,
        font="JetBrains Mono",
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    return s


def bullet_dot(slide, x, y, color=ACCENT_CYAN):
    d = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.12), Inches(0.12))
    d.shadow.inherit = False
    d.line.fill.background()
    d.fill.solid()
    d.fill.fore_color.rgb = color
    return d


# ------------------------------ slides --------------------------------------


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_BLACK)
    grid_lines(slide)

    # decorative: large diagonal stripes
    stripe1 = add_rect(
        slide,
        Inches(-1),
        Inches(2.5),
        Inches(15),
        Emu(11430),  # 0.3in tall
        fill=ACCENT_VIOLET,
        line=False,
    )
    stripe1.rotation = -12
    stripe2 = add_rect(
        slide,
        Inches(-1),
        Inches(4.7),
        Inches(15),
        Emu(8575),
        fill=ACCENT_CYAN,
        line=False,
    )
    stripe2.rotation = -12

    # massive corner ring
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(-2.5), Inches(8), Inches(8))
    ring.shadow.inherit = False
    ring.fill.background()
    ring.line.color.rgb = ACCENT_VIOLET
    ring.line.width = Pt(1.5)

    ring2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.2), Inches(-1.8), Inches(6.6), Inches(6.6))
    ring2.shadow.inherit = False
    ring2.fill.background()
    ring2.line.color.rgb = ACCENT_CYAN
    ring2.line.width = Pt(0.75)

    # eyebrow
    add_text(
        slide,
        "PRIMUS × MEGATRON-LM  ·  PLAN-2 ROADMAP",
        Inches(0.7),
        Inches(0.7),
        Inches(12),
        Inches(0.4),
        size=11,
        color=ACCENT_CYAN,
        bold=True,
        font="JetBrains Mono",
    )

    # title - big two-line
    add_text(
        slide,
        "DEEPSEEK-V4",
        Inches(0.7),
        Inches(2.7),
        Inches(11),
        Inches(1.4),
        size=72,
        color=TEXT,
        bold=True,
        font="Inter",
    )
    add_text(
        slide,
        "in PRIMUS",
        Inches(0.7),
        Inches(3.85),
        Inches(11),
        Inches(1.0),
        size=56,
        color=ACCENT_CYAN,
        bold=True,
        font="Inter",
    )
    add_text(
        slide,
        "Architecture-faithful rewrite. Megatron-native. Checkpoint-compatible.",
        Inches(0.7),
        Inches(4.95),
        Inches(11),
        Inches(0.5),
        size=18,
        color=TEXT_DIM,
        font="Inter",
    )

    # bottom meta line
    add_text(
        slide,
        "v3.0 · 2026-05-01",
        Inches(0.7),
        Inches(6.6),
        Inches(6),
        Inches(0.3),
        size=11,
        color=TEXT_MUTE,
        font="JetBrains Mono",
    )
    add_text(
        slide,
        "PLAN-0  →  PLAN-1  →  PLAN-2",
        Inches(7),
        Inches(6.6),
        Inches(5.6),
        Inches(0.3),
        size=11,
        color=ACCENT_AMBER,
        bold=True,
        font="JetBrains Mono",
        align=PP_ALIGN.RIGHT,
    )

    # progress dots in corner
    for i, c in enumerate([ACCENT_EMERALD, ACCENT_EMERALD, ACCENT_AMBER]):
        d = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(11.6 + i * 0.35),
            Inches(6.55),
            Inches(0.22),
            Inches(0.22),
        )
        d.shadow.inherit = False
        d.line.fill.background()
        d.fill.solid()
        d.fill.fore_color.rgb = c


def slide_overview(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_BLACK)
    grid_lines(slide)
    header_strip(
        slide,
        "01 · executive overview",
        "Three plans, one model, one mission",
        "Where we have been, where we are, and what plan-2 must deliver.",
    )

    # 3 columns: Plan-0 / Plan-1 / Plan-2
    cards = [
        {
            "tag": "PLAN-0",
            "color": ACCENT_AMBER,
            "title": "Bring-up",
            "phases": "P0 → P7",
            "status": "DONE",
            "notes": [
                "Architecture investigation",
                "YAML configs (Flash / Pro)",
                "model_type dispatch in trainer",
                "Initial layer specs / block",
                "HC + hybrid attention scaffold",
                "MoE / activation / RoPE / MTP",
                "Single-node smoke runnable",
            ],
        },
        {
            "tag": "PLAN-1",
            "color": ACCENT_VIOLET,
            "title": "ModuleSpec convergence",
            "phases": "P8 → P11",
            "status": "PARTIAL · PAUSED",
            "notes": [
                "Spec/provider refactor (P9) ✓",
                "TE / Turbo provider reuse (P10) ✓",
                "MoE distributed convergence (P11) –",
                "Hash router still reinvented",
                "Attention still nn.Module",
                "No state-dict adapter",
                "Real V4 parity NOT achieved",
            ],
        },
        {
            "tag": "PLAN-2",
            "color": ACCENT_CYAN,
            "title": "Architecture-faithful rewrite",
            "phases": "P12 → P21",
            "status": "ACTIVE",
            "notes": [
                "Rebase on MLASelfAttention",
                "TransformerLayer / Block reuse",
                "MoELayer / TopKRouter reuse",
                "MultiTokenPredictionBlock reuse",
                "Faithful HC with PP-aware fix",
                "State-dict adapter (HF → primus)",
                "Numerical alignment release gate",
            ],
        },
    ]
    x0 = Inches(0.55)
    y0 = Inches(2.25)
    card_w = Inches(4.0)
    card_h = Inches(4.4)
    gap = Inches(0.13)
    for i, c in enumerate(cards):
        x = x0 + (card_w + gap) * i
        # outer card
        add_rect(slide, x, y0, card_w, card_h, fill=BG_DEEP, line_color=LINE, radius=0.04)
        # accent left bar
        add_rect(slide, x, y0, Inches(0.08), card_h, fill=c["color"], line=False)
        # tag chip
        chip(slide, x + Inches(0.3), y0 + Inches(0.25), c["tag"], color=c["color"], w=Inches(0.95))
        # phases
        add_text(
            slide,
            c["phases"],
            x + Inches(1.4),
            y0 + Inches(0.25),
            Inches(2.5),
            Inches(0.3),
            size=11,
            color=TEXT_MUTE,
            font="JetBrains Mono",
        )
        # title
        add_text(
            slide,
            c["title"],
            x + Inches(0.3),
            y0 + Inches(0.65),
            card_w - Inches(0.3),
            Inches(0.6),
            size=22,
            color=TEXT,
            bold=True,
        )
        # status
        add_text(
            slide,
            c["status"],
            x + Inches(0.3),
            y0 + Inches(1.25),
            card_w - Inches(0.3),
            Inches(0.3),
            size=10,
            color=c["color"],
            bold=True,
            font="JetBrains Mono",
        )
        # divider
        add_line(
            slide,
            x + Inches(0.3),
            y0 + Inches(1.65),
            x + card_w - Inches(0.3),
            y0 + Inches(1.65),
            color=LINE,
        )
        # bullets
        for j, n in enumerate(c["notes"]):
            yy = y0 + Inches(1.85 + j * 0.32)
            bullet_dot(slide, x + Inches(0.35), yy + Inches(0.07), color=c["color"])
            add_text(
                slide,
                n,
                x + Inches(0.6),
                yy,
                card_w - Inches(0.7),
                Inches(0.32),
                size=11.5,
                color=TEXT_DIM,
            )

    # bottom KPI strip
    kpi_y = Inches(6.85)
    kpis = [
        ("Phases shipped", "P0 – P11", ACCENT_AMBER),
        ("Phases re-planned", "P12 – P21", ACCENT_VIOLET),
        ("CRIT findings", "10", ACCENT_RED),
        ("HIGH findings", "11", ACCENT_ROSE),
        ("Reused MCore modules", "+8", ACCENT_EMERALD),
    ]
    kpi_w = Inches(2.4)
    kpi_x0 = Inches(0.55)
    for i, (label, val, c) in enumerate(kpis):
        x = kpi_x0 + (kpi_w + Emu(20000)) * i
        add_text(
            slide,
            val,
            x,
            kpi_y - Inches(0.12),
            kpi_w,
            Inches(0.3),
            size=14,
            color=c,
            bold=True,
            font="JetBrains Mono",
        )
        add_text(
            slide,
            label,
            x + Inches(1.05),
            kpi_y - Inches(0.05),
            kpi_w,
            Inches(0.3),
            size=9,
            color=TEXT_MUTE,
            font="JetBrains Mono",
        )

    footer(slide, page, total)


def slide_v4_architecture(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_BLACK)
    grid_lines(slide)
    header_strip(
        slide,
        "02 · what makes deepseek-v4 hard",
        "The architectural surface we must reproduce",
        "Six axes where V4 diverges from a vanilla MoE transformer.",
    )

    blocks = [
        (
            "Hyper-Connections (HC)",
            "K parallel residual streams with a learnable Sinkhorn-normalised mixer per layer. Splits input into K paths, mixes after each block, collapses at the head.",
            ACCENT_VIOLET,
        ),
        (
            "Hybrid Attention",
            "Mix of Dense / Compressed-Sparse Attn (CSA) / Heavily-Compressed Attn (HCA). Single-latent KV (K=V), per-head q_norm + kv_norm, grouped low-rank O.",
            ACCENT_CYAN,
        ),
        (
            "MoE: Hash + Learned",
            "Hash router for some layers, learnable router with sqrtsoftplus / sigmoid / softmax score for others. Shared experts, expert bias, Megatron MoE dispatch.",
            ACCENT_EMERALD,
        ),
        (
            "Clamped SwiGLU",
            "PRE-multiplication clamp on silu(gate) and up before the elementwise multiply (not post-mul). Distinct w1 / w3 weights — no fused gate-up.",
            ACCENT_AMBER,
        ),
        (
            "Dual-RoPE + YaRN",
            "Two RoPE caches (main, compress) with different theta, partial-RoPE on first dim, YaRN scaling on the compressed path.",
            ACCENT_ROSE,
        ),
        (
            "Multi-Token Prediction",
            "Auxiliary lookahead heads on top of the main decoder. Must reuse Megatron's MultiTokenPredictionBlock for PP / loss / activation checkpointing.",
            ACCENT_CYAN,
        ),
    ]

    x0 = Inches(0.55)
    y0 = Inches(2.25)
    bw = Inches(4.05)
    bh = Inches(2.15)
    gap = Inches(0.15)
    for i, (title, body, color) in enumerate(blocks):
        col = i % 3
        row = i // 3
        x = x0 + (bw + gap) * col
        y = y0 + (bh + gap) * row
        add_rect(slide, x, y, bw, bh, fill=BG_DEEP, line_color=LINE, radius=0.05)
        # number badge
        add_text(
            slide,
            f"{i+1:02d}",
            x + Inches(0.3),
            y + Inches(0.2),
            Inches(0.7),
            Inches(0.5),
            size=24,
            color=color,
            bold=True,
            font="JetBrains Mono",
        )
        # title
        add_text(
            slide,
            title,
            x + Inches(1.2),
            y + Inches(0.25),
            bw - Inches(1.3),
            Inches(0.5),
            size=15,
            color=TEXT,
            bold=True,
        )
        # body
        add_text(
            slide,
            body,
            x + Inches(0.3),
            y + Inches(0.95),
            bw - Inches(0.5),
            bh - Inches(1.05),
            size=11,
            color=TEXT_DIM,
            spacing=1.25,
        )
        # accent corner
        add_rect(slide, x, y, Inches(0.08), bh, fill=color, line=False)

    footer(slide, page, total)


def slide_review_findings(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_BLACK)
    grid_lines(slide)
    header_strip(
        slide,
        "03 · code review findings",
        "What's wrong on dev/wenx/deepseek-v4 today",
        "Severity-ranked findings from the e194e039..HEAD review.",
    )

    # Severity legend strip
    leg_y = Inches(2.15)
    legs = [
        ("CRIT", "Architecture diverges from real V4", ACCENT_RED),
        ("HIGH", "Distributed / spec correctness", ACCENT_ROSE),
        ("MED", "Megatron reuse / hygiene", ACCENT_AMBER),
        ("LOW", "Style / docs / coverage", ACCENT_EMERALD),
    ]
    lx = Inches(0.55)
    for i, (tag, desc, col) in enumerate(legs):
        chip(slide, lx, leg_y, tag, color=col, w=Inches(0.7))
        add_text(
            slide,
            desc,
            lx + Inches(0.8),
            leg_y + Inches(0.02),
            Inches(2.5),
            Inches(0.3),
            size=10,
            color=TEXT_DIM,
        )
        lx += Inches(3.05)

    # Two columns of findings
    findings = [
        (
            "A1",
            "CRIT",
            "Attention uses separate K/V projections (real V4 has single wkv = single-latent MQA).",
            ACCENT_RED,
        ),
        (
            "A2",
            "CRIT",
            "Missing per-head q_norm / kv_norm RMSNorms after the low-rank projections.",
            ACCENT_RED,
        ),
        (
            "A3",
            "CRIT",
            "HashRouter has no learnable gate weight; outputs uniform 1/topk instead of learned scores.",
            ACCENT_RED,
        ),
        (
            "A4",
            "CRIT",
            "Clamped SwiGLU clamps post-mul — real V4 clamps pre-mul on silu(gate) and up.",
            ACCENT_RED,
        ),
        (
            "A5",
            "CRIT",
            "No state-dict adapter — official HF / V4-Flash checkpoints cannot be loaded.",
            ACCENT_RED,
        ),
        (
            "A6",
            "HIGH",
            "HCAAttention RoPE: position_id calculation for compressed layers is off.",
            ACCENT_ROSE,
        ),
        (
            "B1",
            "HIGH",
            "DeepseekV4Attention is plain nn.Module, not a MLASelfAttention subclass.",
            ACCENT_ROSE,
        ),
        ("B2", "HIGH", "DeepseekV4TransformerBlock does not subclass TransformerBlock.", ACCENT_ROSE),
        ("B3", "HIGH", "DeepseekV4HybridLayer does not subclass TransformerLayer.", ACCENT_ROSE),
        ("B4", "HIGH", "DeepseekV4MoE does not subclass MoELayer; reinvents dispatch wiring.", ACCENT_ROSE),
        (
            "C1",
            "HIGH",
            "HC × PP: HyperHead applied per stage, destroys K-stream context across PP boundary.",
            ACCENT_ROSE,
        ),
        (
            "C2",
            "HIGH",
            "Token-IDs stashed on decoder._v4_token_ids — breaks PP, leaks stale state.",
            ACCENT_ROSE,
        ),
        ("C3", "MED", "Position IDs faked internally; caller-supplied position_ids ignored.", ACCENT_AMBER),
        (
            "C4",
            "MED",
            "TP partitioning of attention low-rank uses 'duplicated' instead of sharded.",
            ACCENT_AMBER,
        ),
        ("D1", "MED", "Compressor / Indexer / DualRoPE not expressed as spec submodules.", ACCENT_AMBER),
        (
            "D2",
            "MED",
            "MTP bypasses MultiTokenPredictionBlock; reimplements PP-incompatible variant.",
            ACCENT_AMBER,
        ),
        (
            "E1",
            "LOW",
            "compress_ratios stored as string in YAML, parsed by ast.literal_eval.",
            ACCENT_EMERALD,
        ),
        ("E2", "LOW", "No numerical alignment harness vs HF or NeMo reference.", ACCENT_EMERALD),
    ]

    col_w = Inches(6.1)
    row_h = Inches(0.43)
    y_start = Inches(2.7)
    half = (len(findings) + 1) // 2
    for i, (id_, sev, msg, col) in enumerate(findings):
        col_i = 0 if i < half else 1
        row_i = i if col_i == 0 else (i - half)
        x = Inches(0.55) + (col_w + Inches(0.25)) * col_i
        y = y_start + row_h * row_i
        # row bg
        add_rect(slide, x, y, col_w, Inches(0.38), fill=BG_DEEP, line=False, radius=0.2)
        # accent stripe
        add_rect(slide, x, y, Inches(0.06), Inches(0.38), fill=col, line=False)
        # ID
        add_text(
            slide,
            id_,
            x + Inches(0.18),
            y + Inches(0.05),
            Inches(0.5),
            Inches(0.3),
            size=10,
            color=TEXT_MUTE,
            bold=True,
            font="JetBrains Mono",
        )
        # SEV
        add_text(
            slide,
            sev,
            x + Inches(0.6),
            y + Inches(0.05),
            Inches(0.6),
            Inches(0.3),
            size=10,
            color=col,
            bold=True,
            font="JetBrains Mono",
        )
        # message
        add_text(
            slide,
            msg,
            x + Inches(1.25),
            y + Inches(0.05),
            col_w - Inches(1.35),
            Inches(0.3),
            size=10,
            color=TEXT,
        )

    footer(slide, page, total)


def slide_plan2_strategy(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_BLACK)
    grid_lines(slide)
    header_strip(
        slide,
        "04 · plan-2 strategy",
        "Three pillars: faithful · native · verifiable",
        "Reuse Megatron, parity with reference, gated by numerical tests.",
    )

    pillars = [
        {
            "n": "01",
            "title": "Megatron-Native",
            "color": ACCENT_CYAN,
            "lines": [
                "Subclass MLASelfAttention, TransformerLayer,",
                "TransformerBlock, MoELayer, TopKRouter,",
                "MultiTokenPredictionBlock, YarnRotaryEmbedding.",
                "Express V4 extensions as ModuleSpec submodules:",
                "Compressor / Indexer / DualRoPE / HyperMixer.",
            ],
        },
        {
            "n": "02",
            "title": "Architecture-Faithful",
            "color": ACCENT_VIOLET,
            "lines": [
                "Single-latent KV (wkv), per-head q/kv_norm, ",
                "grouped low-rank O, attention-sink scalar.",
                "HashRouter as a TopKRouter with frozen gate.",
                "Pre-mul clamping for SwiGLU; w1/w3 unfused.",
                "HC fixed across PP via stream-packed send/recv.",
            ],
        },
        {
            "n": "03",
            "title": "Verifiable",
            "color": ACCENT_EMERALD,
            "lines": [
                "State-dict adapter: HF → Primus weight names.",
                "Numerical-alignment harness vs HF reference.",
                "Distributed smokes: TP/PP/EP grids, MTP on/off.",
                "Release gates G1..G13 — every phase contributes.",
                "All gates green before P21 / release.",
            ],
        },
    ]

    x0 = Inches(0.55)
    y0 = Inches(2.3)
    pw = Inches(4.05)
    ph = Inches(4.0)
    gap = Inches(0.13)
    for i, p in enumerate(pillars):
        x = x0 + (pw + gap) * i
        # outer
        add_rect(slide, x, y0, pw, ph, fill=BG_DEEP, line_color=LINE, radius=0.04)
        # huge number
        add_text(
            slide,
            p["n"],
            x + Inches(0.3),
            y0 + Inches(0.25),
            Inches(2),
            Inches(1.4),
            size=72,
            color=p["color"],
            bold=True,
            font="JetBrains Mono",
        )
        # vertical bar
        add_rect(slide, x, y0, Inches(0.08), ph, fill=p["color"], line=False)
        # title
        add_text(
            slide,
            p["title"],
            x + Inches(0.3),
            y0 + Inches(1.55),
            pw - Inches(0.5),
            Inches(0.5),
            size=20,
            color=TEXT,
            bold=True,
        )
        # divider
        add_line(
            slide,
            x + Inches(0.3),
            y0 + Inches(2.05),
            x + pw - Inches(0.3),
            y0 + Inches(2.05),
            color=LINE,
        )
        # bullets
        for j, ln in enumerate(p["lines"]):
            yy = y0 + Inches(2.2 + j * 0.32)
            bullet_dot(slide, x + Inches(0.35), yy + Inches(0.07), color=p["color"])
            add_text(
                slide,
                ln,
                x + Inches(0.6),
                yy,
                pw - Inches(0.7),
                Inches(0.32),
                size=11.5,
                color=TEXT_DIM,
            )

    # bottom guiding-principles ribbon
    rib_y = Inches(6.5)
    add_rect(
        slide, Inches(0.55), rib_y, Inches(12.2), Inches(0.55), fill=BG_PANEL_2, line_color=LINE, radius=0.3
    )
    add_text(
        slide,
        "GUIDING PRINCIPLES  ·  reuse > reinvent  ·  spec > monolith  ·  parity > perf  ·  tests gate phases",
        Inches(0.55),
        rib_y + Inches(0.05),
        Inches(12.2),
        Inches(0.45),
        size=12,
        color=ACCENT_CYAN,
        bold=True,
        font="JetBrains Mono",
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    footer(slide, page, total)


def slide_phases_table(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_BLACK)
    grid_lines(slide)
    header_strip(
        slide,
        "05 · plan-2 phase breakdown",
        "Ten phases — P12 lockdown to P21 release",
        "Each phase has explicit deliverables, exit criteria, and a release gate.",
    )

    phases = [
        ("P12", "Lockdown", "Plan-2 scope sign-off; status.md / dashboards updated.", "DOC", ACCENT_AMBER),
        (
            "P13",
            "Faithful Attention",
            "DeepseekV4Attention(MLASelfAttention) with single-latent KV, q/kv_norm, grouped O, sink scalar.",
            "CODE",
            ACCENT_CYAN,
        ),
        (
            "P14",
            "MoE & Activation",
            "HashRouter as TopKRouter subclass; learned router with sqrtsoftplus; pre-mul clamp; ClampedSwiGLU MLP.",
            "CODE",
            ACCENT_VIOLET,
        ),
        (
            "P15",
            "Layer · Block · HC×PP",
            "DeepseekV4Layer(TransformerLayer); TransformerBlock subclass; HC stream-packed PP send/recv.",
            "CODE",
            ACCENT_VIOLET,
        ),
        (
            "P16",
            "MTP integration",
            "Plug into MultiTokenPredictionBlock; lookahead heads PP-aware; per-head loss aggregation.",
            "CODE",
            ACCENT_CYAN,
        ),
        (
            "P17",
            "State-dict Adapter",
            "HF V4-Flash / V4-Pro → Primus weight name map; round-trip load/save tested.",
            "CODE",
            ACCENT_EMERALD,
        ),
        (
            "P18",
            "Spec System Audit",
            "All V4 modules expressed as ModuleSpec submodules; build_module path everywhere; provider audit.",
            "REFACTOR",
            ACCENT_AMBER,
        ),
        (
            "P19",
            "Distributed Re-validation",
            "TP / PP / EP grid smoke; MTP on/off; HC=K1..K4; FP8 toggles; all gates G1..G11.",
            "TEST",
            ACCENT_ROSE,
        ),
        (
            "P20",
            "Release Gates & Cleanup",
            "Numerical alignment vs HF reference; gate G12 / G13 green; deslop pass.",
            "TEST",
            ACCENT_RED,
        ),
        (
            "P21",
            "Tag & Handoff",
            "Tag v4.0; tech-blog publish; pretrain checkpoint hand-off to convergence team.",
            "RELEASE",
            ACCENT_EMERALD,
        ),
    ]

    head_y = Inches(2.2)
    add_text(
        slide,
        "PHASE",
        Inches(0.55),
        head_y,
        Inches(0.8),
        Inches(0.3),
        size=10,
        color=TEXT_MUTE,
        bold=True,
        font="JetBrains Mono",
    )
    add_text(
        slide,
        "NAME",
        Inches(1.4),
        head_y,
        Inches(2.6),
        Inches(0.3),
        size=10,
        color=TEXT_MUTE,
        bold=True,
        font="JetBrains Mono",
    )
    add_text(
        slide,
        "DELIVERABLE",
        Inches(4.05),
        head_y,
        Inches(7.7),
        Inches(0.3),
        size=10,
        color=TEXT_MUTE,
        bold=True,
        font="JetBrains Mono",
    )
    add_text(
        slide,
        "TYPE",
        Inches(11.85),
        head_y,
        Inches(0.9),
        Inches(0.3),
        size=10,
        color=TEXT_MUTE,
        bold=True,
        font="JetBrains Mono",
        align=PP_ALIGN.RIGHT,
    )
    add_line(
        slide,
        Inches(0.55),
        head_y + Inches(0.32),
        Inches(12.75),
        head_y + Inches(0.32),
        color=LINE_strong if False else LINE,
        weight=1.0,
    )

    row_h = Inches(0.42)
    for i, (p, name, det, kind, color) in enumerate(phases):
        y = head_y + Inches(0.42) + row_h * i
        add_rect(
            slide,
            Inches(0.55),
            y,
            Inches(12.2),
            Inches(0.38),
            fill=BG_DEEP if i % 2 == 0 else BG_BLACK,
            line=False,
            radius=0.1,
        )
        add_rect(slide, Inches(0.55), y, Inches(0.06), Inches(0.38), fill=color, line=False)
        # phase
        add_text(
            slide,
            p,
            Inches(0.7),
            y + Inches(0.05),
            Inches(0.7),
            Inches(0.3),
            size=12,
            color=color,
            bold=True,
            font="JetBrains Mono",
        )
        # name
        add_text(
            slide,
            name,
            Inches(1.4),
            y + Inches(0.05),
            Inches(2.6),
            Inches(0.3),
            size=12,
            color=TEXT,
            bold=True,
        )
        # detail
        add_text(
            slide, det, Inches(4.05), y + Inches(0.05), Inches(7.6), Inches(0.3), size=10.5, color=TEXT_DIM
        )
        # type chip
        chip(slide, Inches(11.85), y + Inches(0.05), kind, color=color, w=Inches(0.85))

    footer(slide, page, total)


# Note: LINE_strong defined for compatibility (unused if False above)
LINE_strong = RGBColor(0x4A, 0x5A, 0x9A)


def slide_dependency_graph(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_BLACK)
    grid_lines(slide)
    header_strip(
        slide,
        "06 · dependency graph",
        "How plan-2 phases unlock each other",
        "Arrows = hard dependencies. Nodes group around three swimlanes.",
    )

    # Swimlanes: Modules / Integration / Validation
    lanes_y = [Inches(2.5), Inches(4.4), Inches(6.0)]
    lanes_label = ["MODULES", "INTEGRATION", "VALIDATION"]
    lanes_color = [ACCENT_CYAN, ACCENT_VIOLET, ACCENT_EMERALD]

    for ly, label, c in zip(lanes_y, lanes_label, lanes_color):
        add_rect(slide, Inches(0.55), ly, Inches(12.2), Inches(0.05), fill=c, line=False)
        add_text(
            slide,
            label,
            Inches(0.55),
            ly - Inches(0.32),
            Inches(2),
            Inches(0.25),
            size=10,
            color=c,
            bold=True,
            font="JetBrains Mono",
        )

    # Node positions: x ∈ [0.7, 12.5]
    nodes = {
        "P12": (Inches(0.85), lanes_y[1] - Inches(0.45), "Lockdown", ACCENT_AMBER, "INTEGRATION"),
        "P13": (Inches(2.4), lanes_y[0] - Inches(0.45), "Faithful Attn", ACCENT_CYAN, "MODULES"),
        "P14": (Inches(4.0), lanes_y[0] - Inches(0.45), "MoE / Activ.", ACCENT_CYAN, "MODULES"),
        "P15": (
            Inches(5.6),
            lanes_y[1] - Inches(0.45),
            "Layer / Block / HC×PP",
            ACCENT_VIOLET,
            "INTEGRATION",
        ),
        "P16": (Inches(7.5), lanes_y[1] - Inches(0.45), "MTP", ACCENT_VIOLET, "INTEGRATION"),
        "P17": (Inches(9.0), lanes_y[1] - Inches(0.45), "State-dict", ACCENT_VIOLET, "INTEGRATION"),
        "P18": (Inches(10.5), lanes_y[1] - Inches(0.45), "Spec audit", ACCENT_AMBER, "INTEGRATION"),
        "P19": (Inches(8.0), lanes_y[2] - Inches(0.45), "Dist. re-valid.", ACCENT_EMERALD, "VALIDATION"),
        "P20": (Inches(10.0), lanes_y[2] - Inches(0.45), "Release gates", ACCENT_RED, "VALIDATION"),
        "P21": (Inches(11.6), lanes_y[2] - Inches(0.45), "Tag", ACCENT_EMERALD, "VALIDATION"),
    }
    nw = Inches(1.4)
    nh = Inches(0.85)

    centers = {}
    for k, (x, y, label, color, lane) in nodes.items():
        add_rect(slide, x, y, nw, nh, fill=BG_DEEP, line_color=color, line_w=1.0, radius=0.15)
        add_text(
            slide,
            k,
            x + Inches(0.1),
            y + Inches(0.07),
            nw,
            Inches(0.3),
            size=11,
            color=color,
            bold=True,
            font="JetBrains Mono",
        )
        add_text(
            slide,
            label,
            x + Inches(0.1),
            y + Inches(0.4),
            nw - Inches(0.2),
            Inches(0.4),
            size=10,
            color=TEXT,
            bold=True,
        )
        centers[k] = (x + nw / 2, y + nh / 2)

    edges = [
        ("P12", "P13"),
        ("P12", "P14"),
        ("P13", "P15"),
        ("P14", "P15"),
        ("P15", "P16"),
        ("P15", "P17"),
        ("P16", "P18"),
        ("P17", "P18"),
        ("P18", "P19"),
        ("P19", "P20"),
        ("P20", "P21"),
    ]
    for a, b in edges:
        ax, ay = centers[a]
        bx, by = centers[b]
        # exit from bottom or right of a, enter top of b
        # use right-mid → left-mid if same lane, else bottom→top
        ax_e = ax + nw / 2 - Inches(0.1) if ay == by else ax
        ay_e = ay if ay == by else ay + nh / 2
        bx_e = bx - nw / 2 + Inches(0.1) if ay == by else bx
        by_e = by if ay == by else by - nh / 2
        # draw
        add_line(slide, ax_e, ay_e, bx_e, by_e, color=LINE_strong, weight=1.0)

    # Legend
    leg_y = Inches(7.0)
    items = [
        ("Module work", ACCENT_CYAN),
        ("Integration / refactor", ACCENT_VIOLET),
        ("Validation / release", ACCENT_EMERALD),
        ("Risk gate", ACCENT_RED),
    ]
    lx = Inches(0.55)
    for label, c in items:
        chip(slide, lx, leg_y, label, color=c, w=Inches(2.1))
        lx += Inches(2.4)

    footer(slide, page, total)


def slide_gantt(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_BLACK)
    grid_lines(slide)
    header_strip(
        slide,
        "08 · timeline",
        "Three plans on one bar chart",
        "Calendar weeks; light bars = done, neon bars = active / planned.",
    )

    # Time axis: weeks W17..W30 (apr → end of jul)
    weeks = ["W17", "W18", "W19", "W20", "W21", "W22", "W23", "W24", "W25", "W26", "W27", "W28", "W29", "W30"]
    chart_x = Inches(2.4)
    chart_y = Inches(2.4)
    chart_w = Inches(10.3)
    chart_h = Inches(4.2)
    n_weeks = len(weeks)
    col_w = chart_w / n_weeks

    # axis line
    add_line(slide, chart_x, chart_y, chart_x + chart_w, chart_y, color=LINE)
    # week ticks
    for i, w in enumerate(weeks):
        x = chart_x + col_w * i
        add_line(slide, x, chart_y, x, chart_y + chart_h, color=RGBColor(0x10, 0x16, 0x28), weight=0.4)
        add_text(
            slide,
            w,
            x,
            chart_y - Inches(0.32),
            col_w,
            Inches(0.25),
            size=9,
            color=TEXT_MUTE,
            font="JetBrains Mono",
            align=PP_ALIGN.CENTER,
        )

    bars = [
        # phase, start_idx, end_idx, color, status
        ("P0  Investigation", 0, 1, ACCENT_AMBER, "done"),
        ("P1  Configs", 0, 1, ACCENT_AMBER, "done"),
        ("P2  model_type", 1, 2, ACCENT_AMBER, "done"),
        ("P3  Layer specs", 1, 2, ACCENT_AMBER, "done"),
        ("P4  HC + Hybrid Attn", 2, 3, ACCENT_AMBER, "done"),
        ("P5  MoE / RoPE / MTP", 2, 3, ACCENT_AMBER, "done"),
        ("P6  Trainer e2e", 3, 4, ACCENT_AMBER, "done"),
        ("P7  Single-node smoke", 3, 4, ACCENT_AMBER, "done"),
        ("P8  Plan-1 scoping", 4, 4, ACCENT_VIOLET, "done"),
        ("P9  Spec refactor", 4, 5, ACCENT_VIOLET, "done"),
        ("P10 TE provider", 5, 6, ACCENT_VIOLET, "done"),
        ("P11 MoE distrib.", 5, 6, ACCENT_VIOLET, "partial"),
        ("P12 Plan-2 lockdown", 6, 7, ACCENT_CYAN, "active"),
        ("P13 Faithful attn", 7, 9, ACCENT_CYAN, "todo"),
        ("P14 MoE & activ.", 8, 10, ACCENT_CYAN, "todo"),
        ("P15 Layer / HC×PP", 9, 11, ACCENT_VIOLET, "todo"),
        ("P16 MTP", 10, 11, ACCENT_VIOLET, "todo"),
        ("P17 State-dict", 10, 12, ACCENT_VIOLET, "todo"),
        ("P18 Spec audit", 11, 12, ACCENT_AMBER, "todo"),
        ("P19 Dist. revalid.", 12, 13, ACCENT_EMERALD, "todo"),
        ("P20 Release gates", 12, 13, ACCENT_RED, "todo"),
        ("P21 Tag & handoff", 13, 14, ACCENT_EMERALD, "todo"),
    ]
    n_rows = len(bars)
    row_h = chart_h / n_rows
    bar_h = row_h * 0.5
    for i, (label, s, e, color, status) in enumerate(bars):
        y = chart_y + row_h * i + (row_h - bar_h) / 2
        # row label
        add_text(
            slide,
            label,
            Inches(0.55),
            y - Inches(0.02),
            chart_x - Inches(0.6),
            bar_h,
            size=8.5,
            color=TEXT_DIM,
            font="Inter",
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # bar
        bx = chart_x + col_w * s
        bw = max(col_w * (e - s), Inches(0.18))
        # Done bars: dimmer fill
        if status == "done":
            fill = BG_PANEL
            line_c = color
        elif status == "partial":
            fill = BG_PANEL_2
            line_c = color
        elif status == "active":
            fill = color
            line_c = color
        else:  # todo
            fill = BG_DEEP
            line_c = color
        rect = add_rect(slide, bx, y, bw, bar_h, fill=fill, line_color=line_c, line_w=1.0, radius=0.4)
        # status badge inside (only if wide enough)
        if bw > Inches(0.6):
            add_text(
                slide,
                status.upper(),
                bx + Inches(0.05),
                y + Inches(0.02),
                bw - Inches(0.1),
                bar_h - Inches(0.02),
                size=7.5,
                color=TEXT if status == "active" else color,
                bold=True,
                font="JetBrains Mono",
                anchor=MSO_ANCHOR.MIDDLE,
            )

    # Legend
    leg_y = Inches(7.0)
    items = [
        ("done", BG_PANEL, ACCENT_AMBER),
        ("partial", BG_PANEL_2, ACCENT_VIOLET),
        ("active", ACCENT_CYAN, ACCENT_CYAN),
        ("planned", BG_DEEP, ACCENT_EMERALD),
    ]
    lx = Inches(0.55)
    for label, fill, line in items:
        add_rect(
            slide, lx, leg_y, Inches(0.35), Inches(0.22), fill=fill, line_color=line, line_w=1.0, radius=0.4
        )
        add_text(
            slide,
            label,
            lx + Inches(0.5),
            leg_y - Inches(0.02),
            Inches(1.4),
            Inches(0.3),
            size=10,
            color=TEXT_DIM,
            font="JetBrains Mono",
        )
        lx += Inches(2.2)

    footer(slide, page, total)


def slide_milestones_risks(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_BLACK)
    grid_lines(slide)
    header_strip(
        slide,
        "09 · milestones & risks",
        "What success looks like — and what could derail it",
        "Stakeholder communication checkpoints + the top-3 things keeping us up at night.",
    )

    # Left: Milestones
    lx = Inches(0.55)
    ly = Inches(2.25)
    lw = Inches(6.1)
    lh = Inches(4.5)
    add_rect(slide, lx, ly, lw, lh, fill=BG_DEEP, line_color=LINE, radius=0.04)
    add_text(
        slide,
        "MILESTONES",
        lx + Inches(0.3),
        ly + Inches(0.25),
        lw - Inches(0.6),
        Inches(0.4),
        size=14,
        color=ACCENT_EMERALD,
        bold=True,
        font="JetBrains Mono",
    )
    add_line(slide, lx + Inches(0.3), ly + Inches(0.7), lx + lw - Inches(0.3), ly + Inches(0.7), color=LINE)

    milestones = [
        ("M1", "Plan-2 sign-off", "End P12", ACCENT_AMBER),
        ("M2", "Faithful attention green", "End P13 — gate G3", ACCENT_CYAN),
        ("M3", "MoE / activation parity", "End P14 — gate G5", ACCENT_VIOLET),
        ("M4", "PP-aware HC verified", "End P15 — gate G7", ACCENT_VIOLET),
        ("M5", "HF checkpoint loadable", "End P17 — gate G9", ACCENT_EMERALD),
        ("M6", "All distributed grids green", "End P19 — gate G11", ACCENT_EMERALD),
        ("M7", "Numerical alignment ≤1e-3", "End P20 — gate G13", ACCENT_RED),
        ("M8", "v4.0 tag & blog published", "End P21", ACCENT_EMERALD),
    ]
    for i, (id_, name, when, col) in enumerate(milestones):
        y = ly + Inches(0.85 + i * 0.42)
        # vertical timeline dot
        bullet_dot(slide, lx + Inches(0.4), y + Inches(0.1), color=col)
        (
            add_line(
                slide,
                lx + Inches(0.46),
                y + Inches(0.22),
                lx + Inches(0.46),
                y + Inches(0.42 + 0.05),
                color=LINE,
                weight=0.5,
            )
            if i < len(milestones) - 1
            else None
        )
        add_text(
            slide,
            id_,
            lx + Inches(0.7),
            y + Inches(0.04),
            Inches(0.5),
            Inches(0.3),
            size=11,
            color=col,
            bold=True,
            font="JetBrains Mono",
        )
        add_text(
            slide,
            name,
            lx + Inches(1.2),
            y + Inches(0.04),
            Inches(3.0),
            Inches(0.3),
            size=11.5,
            color=TEXT,
            bold=True,
        )
        add_text(
            slide,
            when,
            lx + Inches(4.2),
            y + Inches(0.04),
            Inches(2.0),
            Inches(0.3),
            size=10,
            color=TEXT_DIM,
            font="JetBrains Mono",
            align=PP_ALIGN.RIGHT,
        )

    # Right: Top risks
    rx = Inches(6.85)
    ry = ly
    rw = Inches(5.9)
    rh = lh
    add_rect(slide, rx, ry, rw, rh, fill=BG_DEEP, line_color=LINE, radius=0.04)
    add_text(
        slide,
        "TOP RISKS",
        rx + Inches(0.3),
        ry + Inches(0.25),
        rw - Inches(0.6),
        Inches(0.4),
        size=14,
        color=ACCENT_RED,
        bold=True,
        font="JetBrains Mono",
    )
    add_line(slide, rx + Inches(0.3), ry + Inches(0.7), rx + rw - Inches(0.3), ry + Inches(0.7), color=LINE)

    risks = [
        (
            "R1",
            "Real V4 reference drift",
            "HF reference and NeMo diverge on partial-RoPE pairing and HC head application; we may need two compatibility modes.",
            ACCENT_RED,
        ),
        (
            "R2",
            "HC × PP correctness",
            "K-stream packing across PP boundaries hasn't been done in MCore; needs custom send/recv tensor schema.",
            ACCENT_ROSE,
        ),
        (
            "R3",
            "State-dict ABI",
            "Per-head q_norm / kv_norm naming differs between HF and NeMo; must lock canonical naming before P17.",
            ACCENT_AMBER,
        ),
        (
            "R4",
            "FP8 / Turbo regressions",
            "Re-rooting attention on MLASelfAttention may break TE / Turbo quantised paths. Needs explicit P19 grid.",
            ACCENT_AMBER,
        ),
        (
            "R5",
            "Numerical-alignment thresholds",
            "Mixing TE FP8 + ClampedSwiGLU + sinkhorn HC produces noise > 1e-3. Need calibrated thresholds per gate.",
            ACCENT_VIOLET,
        ),
    ]
    for i, (id_, name, body, col) in enumerate(risks):
        y = ry + Inches(0.85 + i * 0.72)
        chip(slide, rx + Inches(0.3), y, id_, color=col, w=Inches(0.55))
        add_text(
            slide,
            name,
            rx + Inches(0.95),
            y - Inches(0.02),
            rw - Inches(1.2),
            Inches(0.3),
            size=12,
            color=TEXT,
            bold=True,
        )
        add_text(
            slide,
            body,
            rx + Inches(0.95),
            y + Inches(0.28),
            rw - Inches(1.2),
            Inches(0.5),
            size=10,
            color=TEXT_DIM,
            spacing=1.2,
        )

    footer(slide, page, total)


def slide_test_strategy(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_BLACK)
    grid_lines(slide)
    header_strip(
        slide,
        "10 · test strategy",
        "Four levels, thirteen gates, one release",
        "Each phase contributes tests; nothing ships without G1..G13 green.",
    )

    # 4 level columns
    levels = [
        {
            "tag": "L0",
            "title": "Unit",
            "color": ACCENT_CYAN,
            "items": [
                "RoPE / RMSNorm exactness",
                "ClampedSwiGLU pre-mul clamp",
                "TopKRouter scores",
                "HashRouter tid2eid stability",
            ],
        },
        {
            "tag": "L1",
            "title": "Module integration",
            "color": ACCENT_VIOLET,
            "items": [
                "MLA single-latent KV",
                "Compressor / Indexer wiring",
                "DualRoPE on attention block",
                "MoE dispatcher round-trip",
            ],
        },
        {
            "tag": "L2",
            "title": "Distributed smoke",
            "color": ACCENT_AMBER,
            "items": [
                "TP=2/4 ; PP=2/4 ; EP=4/8",
                "MTP on/off",
                "HC K=1..4",
                "Activation checkpoint on/off",
                "FP8 Turbo on/off",
            ],
        },
        {
            "tag": "L3",
            "title": "Release gate",
            "color": ACCENT_EMERALD,
            "items": [
                "HF reference numerical align ≤ 1e-3",
                "Checkpoint round-trip",
                "Loss curve sanity (toy)",
                "Throughput vs Plan-1 ≥ 0.95×",
            ],
        },
    ]
    x0 = Inches(0.55)
    y0 = Inches(2.25)
    cw = Inches(2.99)
    ch = Inches(3.6)
    for i, lv in enumerate(levels):
        x = x0 + (cw + Inches(0.1)) * i
        add_rect(slide, x, y0, cw, ch, fill=BG_DEEP, line_color=LINE, radius=0.04)
        add_rect(slide, x, y0, cw, Inches(0.7), fill=lv["color"], line=False, radius=0.04)
        add_text(
            slide,
            lv["tag"],
            x + Inches(0.25),
            y0 + Inches(0.13),
            Inches(0.7),
            Inches(0.4),
            size=18,
            color=BG_BLACK,
            bold=True,
            font="JetBrains Mono",
        )
        add_text(
            slide,
            lv["title"],
            x + Inches(1.0),
            y0 + Inches(0.18),
            cw - Inches(1.1),
            Inches(0.4),
            size=14,
            color=BG_BLACK,
            bold=True,
        )
        for j, it in enumerate(lv["items"]):
            yy = y0 + Inches(0.95 + j * 0.42)
            bullet_dot(slide, x + Inches(0.25), yy + Inches(0.07), color=lv["color"])
            add_text(
                slide,
                it,
                x + Inches(0.5),
                yy,
                cw - Inches(0.7),
                Inches(0.4),
                size=11,
                color=TEXT_DIM,
                spacing=1.2,
            )

    # Gate strip at bottom
    gate_y = Inches(6.1)
    add_text(
        slide,
        "RELEASE GATES",
        Inches(0.55),
        gate_y,
        Inches(3),
        Inches(0.3),
        size=11,
        color=ACCENT_EMERALD,
        bold=True,
        font="JetBrains Mono",
    )
    gates = [
        "G1 build",
        "G2 lint",
        "G3 attn",
        "G4 csa/hca",
        "G5 moe",
        "G6 swiglu",
        "G7 hc-pp",
        "G8 mtp",
        "G9 ckpt",
        "G10 spec",
        "G11 dist",
        "G12 align",
        "G13 release",
    ]
    gx = Inches(0.55)
    gy = gate_y + Inches(0.4)
    gw = Inches(0.92)
    for i, g in enumerate(gates):
        x = gx + (gw + Emu(20000)) * i
        add_rect(
            slide, x, gy, gw, Inches(0.35), fill=None, line_color=ACCENT_EMERALD, line_w=0.75, radius=0.4
        )
        add_text(
            slide,
            g,
            x,
            gy + Inches(0.04),
            gw,
            Inches(0.3),
            size=9,
            color=ACCENT_EMERALD,
            bold=True,
            font="JetBrains Mono",
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    footer(slide, page, total)


def slide_dev_schedule(prs, page, total):
    """3-row day-by-day schedule: date / phase A~B / work content + flow arrow.

    Past phases (Apr 28 → May 01) and 4 planned days (May 06 → May 09).
    May 02 → 05 is a national holiday and is NOT shown as a column.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_BLACK)
    grid_lines(slide)
    header_strip(
        slide,
        "07 · 开发计划  ·  development schedule",
        "Day-by-day plan — 8 working days across two work blocks",
        "Past block: Apr 28 → May 01 (landed). Holiday: May 02 → 05 (skipped). Plan-2 block: May 06 → 09.",
    )

    # Eight day columns. Layout:
    #   row 1 (h1): date chip
    #   row 2 (h2): phase A~B chip
    #   row 3 (h3): work content card (taller)
    #   horizontal arrow runs through the middle of row 1
    days = [
        # done block
        {
            "date": "Apr 28",
            "weekday": "Tue",
            "phase": "P0 ~ P7",
            "color": ACCENT_EMERALD,
            "status": "DONE",
            "work_title": "MVP bring-up",
            "work_lines": [
                "Investigation + tech blog",
                "YAML configs (base / Flash / Pro)",
                "model_type dispatch + builder",
                "HC + Hybrid Attn + MoE/MTP",
                "1×8 PP=2 EP=4 BF16 smoke",
            ],
        },
        {
            "date": "Apr 29",
            "weekday": "Wed",
            "phase": "P8",
            "color": ACCENT_EMERALD,
            "status": "DONE",
            "work_title": "ModuleSpec refactor",
            "work_lines": [
                "Spec-driven runtime path",
                "DeepseekV4Model on LanguageModule",
                "Builder routes V4 spec only",
                "Retire GPT-placeholder swap",
            ],
        },
        {
            "date": "Apr 30",
            "weekday": "Thu",
            "phase": "P9 ~ P10",
            "color": ACCENT_EMERALD,
            "status": "DONE",
            "work_title": "TE provider + MoE bridge",
            "work_lines": [
                "DeepSeekV4SpecProvider",
                "Norm + projection providerized",
                "MoE submodules + dispatcher",
                "EP all-reduce fallback gated",
                "A/B local↔TE forward report",
            ],
        },
        {
            "date": "May 01",
            "weekday": "Fri",
            "phase": "P12",
            "color": ACCENT_CYAN,
            "status": "DONE",
            "work_title": "Plan-2 lockdown",
            "work_lines": [
                "Architecture review",
                "10 CRIT / 11 HIGH findings",
                "plan-2/00..04 + README",
                "status.md P12+ section",
                "timeline.html + PPT roadmap",
            ],
        },
        # planned block
        {
            "date": "May 06",
            "weekday": "Tue",
            "phase": "P13 ~ P14",
            "color": ACCENT_VIOLET,
            "status": "PLAN",
            "work_title": "Faithful attn + MoE",
            "work_lines": [
                "DeepseekV4Attention(MLA)",
                "Single-latent KV (wkv)",
                "q_norm / kv_norm / sink",
                "Grouped low-rank O proj",
                "ClampedSwiGLU pre-mul",
                "Hash + learned routers",
            ],
        },
        {
            "date": "May 07",
            "weekday": "Wed",
            "phase": "P15 ~ P16",
            "color": ACCENT_VIOLET,
            "status": "PLAN",
            "work_title": "Layer/Block/HC×PP + MTP",
            "work_lines": [
                "DeepseekV4HybridLayer(TL)",
                "DeepseekV4Block(TB)",
                "[S,B,K,D] PP packing",
                "HyperHead on post_process",
                "MultiTokenPredictionBlock",
                "process_mtp_loss wired",
            ],
        },
        {
            "date": "May 08",
            "weekday": "Thu",
            "phase": "P17 ~ P18",
            "color": ACCENT_VIOLET,
            "status": "PLAN",
            "work_title": "Checkpoint + spec audit",
            "work_lines": [
                "HF→Primus state-dict map",
                "load_v4_flash_check.py",
                "tid2eid as buffer",
                "Round-trip + token-0 align",
                "Provider per build call",
                "compress_ratios → list",
            ],
        },
        {
            "date": "May 09",
            "weekday": "Fri",
            "phase": "P19 ~ P21",
            "color": ACCENT_AMBER,
            "status": "PLAN",
            "work_title": "Validation + release",
            "work_lines": [
                "TP/PP/EP smoke grid",
                "Routing snapshot diff = 0",
                "G9 numerical alignment",
                "200-step convergence vs HF",
                "TE on/off perf report",
                "Cleanup + tag + handover",
            ],
        },
    ]

    n = len(days)
    margin_x = Inches(0.45)
    content_w = SLIDE_W - 2 * margin_x
    gap = Inches(0.08)
    col_w = (content_w - gap * (n - 1)) / n

    # Y layout
    y_dates = Inches(2.30)  # row 1: date chips
    h_dates = Inches(0.55)
    y_arrow_mid = y_dates + h_dates + Inches(0.18)  # arrow between row 1 and row 2

    y_phase = y_dates + h_dates + Inches(0.45)  # row 2: phase chips
    h_phase = Inches(0.45)

    y_work = y_phase + h_phase + Inches(0.20)  # row 3: work cards
    h_work = Inches(3.10)

    # Find boundary between done and plan blocks
    last_done_idx = max(i for i, d in enumerate(days) if d["status"] == "DONE")
    first_plan_idx = last_done_idx + 1

    # Holiday gap visualisation: between last_done_idx and first_plan_idx
    boundary_x = margin_x + (col_w + gap) * first_plan_idx - gap / 2

    # ---- horizontal flow arrow (row 1 mid) ----
    arrow_y = y_arrow_mid
    arrow_start_x = margin_x + Inches(0.05)
    arrow_end_x = margin_x + content_w - Inches(0.05)

    # Pre-holiday arrow (emerald, solid line)
    seg1_end_x = boundary_x - Inches(0.35)
    add_line(slide, arrow_start_x, arrow_y, seg1_end_x, arrow_y, color=ACCENT_EMERALD, weight=2.5)

    # Holiday gap (rose dashed indicator with label)
    holiday_w = Inches(0.7)
    hx = boundary_x - holiday_w / 2
    add_rect(
        slide,
        hx,
        arrow_y - Inches(0.18),
        holiday_w,
        Inches(0.36),
        fill=BG_DEEP,
        line_color=ACCENT_ROSE,
        line_w=1.0,
        radius=0.4,
    )
    add_text(
        slide,
        "May 2-5\nHOLIDAY",
        hx,
        arrow_y - Inches(0.18),
        holiday_w,
        Inches(0.36),
        size=8,
        color=ACCENT_ROSE,
        bold=True,
        font="Consolas",
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
        spacing=1.0,
    )

    # Post-holiday arrow (violet, with arrowhead)
    seg2_start_x = boundary_x + Inches(0.35)
    arrow_shaft = slide.shapes.add_connector(1, seg2_start_x, arrow_y, arrow_end_x, arrow_y)
    arrow_shaft.shadow.inherit = False
    arrow_shaft.line.color.rgb = ACCENT_VIOLET
    arrow_shaft.line.width = Pt(2.5)
    # Arrowhead at end (use a simple right-triangle)
    head = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        arrow_end_x - Inches(0.05),
        arrow_y - Inches(0.12),
        Inches(0.18),
        Inches(0.24),
    )
    head.shadow.inherit = False
    head.fill.solid()
    head.fill.fore_color.rgb = ACCENT_VIOLET
    head.line.fill.background()
    head.rotation = 90  # rotate so the right angle points right

    # ---- per-column content ----
    for i, d in enumerate(days):
        x = margin_x + (col_w + gap) * i

        # ROW 1: date chip
        is_done = d["status"] == "DONE"
        date_fill = BG_DEEP if is_done else BG_PANEL
        date_line = ACCENT_EMERALD if is_done else d["color"]
        add_rect(
            slide, x, y_dates, col_w, h_dates, fill=date_fill, line_color=date_line, line_w=1.25, radius=0.15
        )
        # date string
        add_text(
            slide,
            d["date"],
            x,
            y_dates + Inches(0.05),
            col_w,
            Inches(0.30),
            size=14,
            color=date_line,
            bold=True,
            font="Consolas",
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text(
            slide,
            d["weekday"] + "  ·  " + d["status"],
            x,
            y_dates + Inches(0.30),
            col_w,
            Inches(0.22),
            size=8,
            color=TEXT_MUTE,
            font="Consolas",
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # date marker dot on the arrow line
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x + col_w / 2 - Inches(0.07),
            arrow_y - Inches(0.07),
            Inches(0.14),
            Inches(0.14),
        )
        dot.shadow.inherit = False
        dot.line.fill.background()
        dot.fill.solid()
        dot.fill.fore_color.rgb = date_line

        # ROW 2: phase chip (small, rounded)
        add_rect(slide, x, y_phase, col_w, h_phase, fill=d["color"], line=False, radius=0.4)
        add_text(
            slide,
            d["phase"],
            x,
            y_phase,
            col_w,
            h_phase,
            size=14,
            color=BG_BLACK,
            bold=True,
            font="Consolas",
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # ROW 3: work content card (taller)
        add_rect(slide, x, y_work, col_w, h_work, fill=BG_DEEP, line_color=LINE, radius=0.05)
        # accent stripe at top of card
        add_rect(slide, x, y_work, col_w, Inches(0.06), fill=d["color"], line=False, radius=0.05)
        # work title
        add_text(
            slide,
            d["work_title"],
            x + Inches(0.12),
            y_work + Inches(0.18),
            col_w - Inches(0.24),
            Inches(0.55),
            size=11,
            color=TEXT,
            bold=True,
        )
        # divider
        add_line(
            slide,
            x + Inches(0.12),
            y_work + Inches(0.65),
            x + col_w - Inches(0.12),
            y_work + Inches(0.65),
            color=LINE,
        )
        # bullet lines
        for j, ln in enumerate(d["work_lines"]):
            yy = y_work + Inches(0.78 + j * 0.34)
            bullet_dot(slide, x + Inches(0.18), yy + Inches(0.06), color=d["color"])
            add_text(
                slide,
                ln,
                x + Inches(0.36),
                yy,
                col_w - Inches(0.5),
                Inches(0.32),
                size=9.5,
                color=TEXT_DIM,
                spacing=1.15,
            )

    # Block labels at top of two halves
    add_text(
        slide,
        "BLOCK A · Apr 28 → May 01  ·  4 days  ·  landed",
        margin_x,
        Inches(2.00),
        col_w * (last_done_idx + 1) + gap * last_done_idx,
        Inches(0.25),
        size=9,
        color=ACCENT_EMERALD,
        bold=True,
        font="Consolas",
        align=PP_ALIGN.LEFT,
    )
    add_text(
        slide,
        "BLOCK B · May 06 → May 09  ·  4 days  ·  plan-2 execution",
        margin_x + (col_w + gap) * first_plan_idx,
        Inches(2.00),
        col_w * (n - first_plan_idx) + gap * (n - first_plan_idx - 1),
        Inches(0.25),
        size=9,
        color=ACCENT_VIOLET,
        bold=True,
        font="Consolas",
        align=PP_ALIGN.LEFT,
    )

    footer(slide, page, total)


def slide_status_next(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_BLACK)
    grid_lines(slide)
    header_strip(
        slide,
        "11 · current status & next-2-weeks",
        "Where we are today and what lands next",
        "Snapshot taken 2026-05-01 against branch dev/wenx/deepseek-v4.",
    )

    # Left: today
    lx = Inches(0.55)
    ly = Inches(2.25)
    lw = Inches(6.1)
    lh = Inches(4.5)
    add_rect(slide, lx, ly, lw, lh, fill=BG_DEEP, line_color=LINE, radius=0.04)
    add_text(
        slide,
        "TODAY  ·  PLAN-2 LOCKDOWN",
        lx + Inches(0.3),
        ly + Inches(0.25),
        lw - Inches(0.6),
        Inches(0.4),
        size=13,
        color=ACCENT_CYAN,
        bold=True,
        font="JetBrains Mono",
    )
    add_line(slide, lx + Inches(0.3), ly + Inches(0.7), lx + lw - Inches(0.3), ly + Inches(0.7), color=LINE)

    today_items = [
        ("✓", "Architecture review done", "10 CRIT / 11 HIGH / 6 MED / 5 LOW", ACCENT_EMERALD),
        ("✓", "Plan-2 docs landed", "README + 00..04 in plan-2/", ACCENT_EMERALD),
        ("✓", "status.md updated", "P12 lockdown row appended", ACCENT_EMERALD),
        ("✓", "Timeline HTML", "progress/timeline.html", ACCENT_EMERALD),
        ("⟶", "Awaiting sign-off", "stakeholder review of plan-2 scope", ACCENT_AMBER),
        ("⟶", "Tech-blog rollover", "as-built notes for plan-1 + pointer to plan-2", ACCENT_AMBER),
    ]
    for i, (mark, title, body, col) in enumerate(today_items):
        y = ly + Inches(0.9 + i * 0.6)
        add_text(
            slide,
            mark,
            lx + Inches(0.35),
            y - Inches(0.02),
            Inches(0.4),
            Inches(0.4),
            size=18,
            color=col,
            bold=True,
            font="JetBrains Mono",
        )
        add_text(
            slide, title, lx + Inches(0.85), y, lw - Inches(1.1), Inches(0.3), size=12, color=TEXT, bold=True
        )
        add_text(
            slide,
            body,
            lx + Inches(0.85),
            y + Inches(0.27),
            lw - Inches(1.1),
            Inches(0.3),
            size=10,
            color=TEXT_DIM,
        )

    # Right: next-2-weeks
    rx = Inches(6.85)
    ry = ly
    rw = Inches(5.9)
    rh = lh
    add_rect(slide, rx, ry, rw, rh, fill=BG_DEEP, line_color=LINE, radius=0.04)
    add_text(
        slide,
        "NEXT 4 WORKDAYS  ·  MAY 06 → MAY 09",
        rx + Inches(0.3),
        ry + Inches(0.25),
        rw - Inches(0.6),
        Inches(0.4),
        size=13,
        color=ACCENT_VIOLET,
        bold=True,
        font="Consolas",
    )
    add_line(slide, rx + Inches(0.3), ry + Inches(0.7), rx + rw - Inches(0.3), ry + Inches(0.7), color=LINE)

    next_items = [
        (
            "May 06",
            "P13 + P14  ·  Faithful attn + MoE",
            "MLA-rooted attn with single-latent KV / q+kv norm / grouped O / sink. Hash + learned routers; pre-mul ClampedSwiGLU.",
            ACCENT_CYAN,
        ),
        (
            "May 07",
            "P15 + P16  ·  Layer/Block + MTP",
            "DeepseekV4HybridLayer(TL) + Block(TB) with [S,B,K,D] PP packing. MultiTokenPredictionBlock + per-MTP HyperHead.",
            ACCENT_VIOLET,
        ),
        (
            "May 08",
            "P17 + P18  ·  Checkpoint + spec audit",
            "HF→Primus state-dict adapter; load V4-Flash; provider-built-once contract; YAML compress_ratios → list.",
            ACCENT_VIOLET,
        ),
        (
            "May 09",
            "P19 + P20 + P21  ·  Validate + ship",
            "TP/PP/EP grid smoke; 200-step convergence vs HF; TE on/off perf report; cleanup + tag + handover.",
            ACCENT_AMBER,
        ),
    ]
    for i, (when, title, body, col) in enumerate(next_items):
        y = ry + Inches(0.9 + i * 0.7)
        chip(slide, rx + Inches(0.3), y, when, color=col, w=Inches(1.2))
        add_text(
            slide,
            title,
            rx + Inches(1.6),
            y - Inches(0.02),
            rw - Inches(1.85),
            Inches(0.3),
            size=12,
            color=TEXT,
            bold=True,
        )
        add_text(
            slide,
            body,
            rx + Inches(1.6),
            y + Inches(0.27),
            rw - Inches(1.85),
            Inches(0.5),
            size=10,
            color=TEXT_DIM,
            spacing=1.2,
        )

    footer(slide, page, total)


def slide_thanks(prs, page, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_BLACK)
    grid_lines(slide)

    # giant ring
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-3), Inches(-3), Inches(13), Inches(13))
    ring.shadow.inherit = False
    ring.fill.background()
    ring.line.color.rgb = ACCENT_VIOLET
    ring.line.width = Pt(1.0)
    ring2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(-1.5), Inches(10), Inches(10))
    ring2.shadow.inherit = False
    ring2.fill.background()
    ring2.line.color.rgb = ACCENT_CYAN
    ring2.line.width = Pt(0.5)

    # title
    add_text(
        slide,
        "READY FOR PLAN-2",
        Inches(5.2),
        Inches(2.8),
        Inches(8),
        Inches(0.4),
        size=12,
        color=ACCENT_CYAN,
        bold=True,
        font="JetBrains Mono",
    )
    add_text(
        slide, "Let's build", Inches(5.2), Inches(3.2), Inches(8), Inches(1.0), size=44, color=TEXT, bold=True
    )
    add_text(
        slide,
        "a real DeepSeek-V4.",
        Inches(5.2),
        Inches(4.05),
        Inches(8),
        Inches(1.0),
        size=44,
        color=ACCENT_CYAN,
        bold=True,
    )

    add_text(
        slide,
        "Documents:  deepseek-v4/develop/plan-2/  ·  Progress:  deepseek-v4/develop/progress/",
        Inches(0.55),
        Inches(7.05),
        Inches(12),
        Inches(0.3),
        size=10,
        color=TEXT_MUTE,
        font="JetBrains Mono",
    )
    add_text(
        slide,
        f"{page:02d} / {total:02d}",
        Inches(11.8),
        Inches(7.05),
        Inches(1.0),
        Inches(0.3),
        size=9,
        color=TEXT_MUTE,
        font="JetBrains Mono",
        align=PP_ALIGN.RIGHT,
    )


# ------------------------------ build ---------------------------------------


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Build all slides
    builders = [
        ("cover", slide_cover),
        ("overview", slide_overview),
        ("v4_arch", slide_v4_architecture),
        ("findings", slide_review_findings),
        ("strategy", slide_plan2_strategy),
        ("phases", slide_phases_table),
        ("dep_graph", slide_dependency_graph),
        ("dev_schedule", slide_dev_schedule),
        ("gantt", slide_gantt),
        ("milestones", slide_milestones_risks),
        ("tests", slide_test_strategy),
        ("status", slide_status_next),
        ("thanks", slide_thanks),
    ]
    total = len(builders)
    for i, (name, fn) in enumerate(builders, start=1):
        if name == "cover":
            fn(prs)
        else:
            fn(prs, i, total)

    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH}")
    print(f"Slides: {total}")


if __name__ == "__main__":
    main()
